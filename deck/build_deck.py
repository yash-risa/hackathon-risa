#!/usr/bin/env python3
"""Build the Auth Guardian hackathon slide deck (editable .pptx).

All facts/numbers are sourced from research-brief.md and RESULTS.md.
Run: python3 build_deck.py  ->  writes auth-guardian.pptx next to this file.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---------------------------------------------------------------- theme
ACCENT      = RGBColor(0x0E, 0x74, 0x90)   # teal
ACCENT_DK   = RGBColor(0x0A, 0x4D, 0x60)   # deep teal
ACCENT_SOFT = RGBColor(0xE2, 0xF1, 0xF4)   # pale teal panel
INK         = RGBColor(0x14, 0x23, 0x3A)   # near-black slate
BODY        = RGBColor(0x33, 0x42, 0x55)   # body gray
MUTED       = RGBColor(0x6B, 0x7A, 0x8C)   # muted gray
LINE        = RGBColor(0xD8, 0xDF, 0xE6)   # hairline
PANEL       = RGBColor(0xF5, 0xF8, 0xFA)   # light panel
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)   # reactive / risk
GREEN       = RGBColor(0x1E, 0x8E, 0x5A)   # good
AMBER       = RGBColor(0xC8, 0x8A, 0x1A)

FONT = "Calibri"
FONT_H = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def _no_line(shape):
    shape.line.fill.background()


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def setpara(p, text, size, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT,
            space_after=6, space_before=0, line=1.05, italic=False):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    try:
        p.line_spacing = line
    except Exception:
        pass
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return p


def add_para(tf, text, size, **kw):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs
                             and not tf.paragraphs[0].text) else tf.add_paragraph()
    return setpara(p, text, size, **kw)


def bullets(tf, items, size=18, color=BODY, gap=10, line=1.05, bold_lead=False):
    first = True
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
        else:
            lead, rest = None, it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        try:
            p.line_spacing = line
        except Exception:
            pass
        # bullet glyph run
        rb = p.add_run()
        rb.text = "\u2022  "
        rb.font.size = Pt(size)
        rb.font.name = FONT
        rb.font.color.rgb = ACCENT
        rb.font.bold = True
        if lead:
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.name = FONT
            r1.font.bold = True
            r1.font.color.rgb = INK
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.name = FONT
            r2.font.color.rgb = color
        else:
            r1 = p.add_run()
            r1.text = rest
            r1.font.size = Pt(size)
            r1.font.name = FONT
            r1.font.bold = bold_lead
            r1.font.color.rgb = color
    return tf


def kicker_title(s, kicker, title, num):
    # left accent bar
    rect(s, Inches(0), Inches(0), Inches(0.18), EMU_H, fill=ACCENT)
    tb, tf = textbox(s, Inches(0.7), Inches(0.5), Inches(11.4), Inches(1.4))
    add_para(tf, kicker.upper(), 13, color=ACCENT, bold=True, space_after=4)
    p = tf.paragraphs[-1]
    for r in p.runs:
        r.font.name = FONT_H
    # letter spacing-ish via spaces not needed
    add_para(tf, title, 30, color=INK, bold=True, line=1.0)
    for r in tf.paragraphs[-1].runs:
        r.font.name = FONT_H
    # slide number
    tbn, tfn = textbox(s, Inches(12.4), Inches(6.95), Inches(0.7), Inches(0.4))
    setpara(tfn.paragraphs[0], num, 11, color=MUTED, align=PP_ALIGN.RIGHT, space_after=0)
    # footer hairline
    rect(s, Inches(0.7), Inches(6.92), Inches(11.4), Pt(1), fill=LINE)
    return


def footnote(s, text, color=MUTED):
    tb, tf = textbox(s, Inches(0.7), Inches(6.98), Inches(11.2), Inches(0.4))
    setpara(tf.paragraphs[0], text, 9.5, color=color, italic=True, space_after=0)


def stat_card(s, x, y, w, h, big, label, sub=None, accent=ACCENT, big_color=None):
    card = rect(s, x, y, w, h, fill=WHITE, line=LINE, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _round(card)
    rect(s, x, y, Inches(0.09), h, fill=accent)  # accent edge
    tb, tf = textbox(s, x + Inches(0.28), y + Inches(0.16), w - Inches(0.4), h - Inches(0.3),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, big, 30, color=big_color or accent, bold=True, space_after=2, line=0.95)
    for r in tf.paragraphs[-1].runs:
        r.font.name = FONT_H
    add_para(tf, label, 12.5, color=INK, bold=True, space_after=1, line=1.0)
    if sub:
        add_para(tf, sub, 10.5, color=MUTED, space_after=0, line=1.0)
    return card


def _round(shape, radius=0.08):
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass


def panel(s, x, y, w, h, fill=PANEL, line=LINE):
    p = rect(s, x, y, w, h, fill=fill, line=line, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _round(p, 0.04)
    return p


# ================================================================ SLIDE 1 — TITLE
s = slide()
# full bleed background
bg = rect(s, 0, 0, EMU_W, EMU_H, fill=INK)
# accent side band
rect(s, 0, 0, Inches(0.35), EMU_H, fill=ACCENT)
# subtle deep panel block
rect(s, Inches(8.7), 0, Inches(4.633), EMU_H, fill=ACCENT_DK)

tb, tf = textbox(s, Inches(0.95), Inches(1.45), Inches(7.6), Inches(0.5))
setpara(tf.paragraphs[0], "RISA LABS  ·  PROACTIVE PA-CONTINUITY LAYER", 13,
        color=RGBColor(0x8E, 0xC7, 0xD4), bold=True, space_after=0)

tb, tf = textbox(s, Inches(0.92), Inches(2.1), Inches(8.0), Inches(2.2))
setpara(tf.paragraphs[0], "Auth Guardian", 66, color=WHITE, bold=True, space_after=6, line=0.98)
tf.paragraphs[0].runs[0].font.name = FONT_H

tb, tf = textbox(s, Inches(0.95), Inches(3.7), Inches(7.3), Inches(1.3))
setpara(tf.paragraphs[0],
        "Keeping oncology authorizations alive \u2014 before the therapy gap happens.",
        24, color=RGBColor(0xDD, 0xEA, 0xEF), bold=False, line=1.1, space_after=0)

# subtitle / descriptor
tb, tf = textbox(s, Inches(0.95), Inches(5.25), Inches(7.2), Inches(1.0))
setpara(tf.paragraphs[0],
        "A Risa Labs proactive PA-continuity layer that detects renewal risk early, "
        "waits for client approval, then files through Risa\u2019s existing PA engine.",
        14.5, color=RGBColor(0xA9, 0xBC, 0xC8), line=1.25, space_after=0)

# right band callouts
rb, rtf = textbox(s, Inches(9.05), Inches(1.7), Inches(4.0), Inches(4.2))
for i, (big, lab) in enumerate([
        ("0", "lapses in backtest (down from 47)"),
        ("$470K", "claim value protected"),
        ("2,024", "therapy-gap days eliminated"),
        ("Zero", "new integrations required")]):
    p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
    setpara(p, big, 34, color=WHITE, bold=True, space_after=0, space_before=(0 if i == 0 else 16), line=0.95)
    p.runs[0].font.name = FONT_H
    p2 = rtf.add_paragraph()
    setpara(p2, lab, 12, color=RGBColor(0xB9, 0xD7, 0xDF), space_after=0, line=1.05)

foottb, foottf = textbox(s, Inches(0.95), Inches(6.95), Inches(7.5), Inches(0.4))
setpara(foottf.paragraphs[0],
        "Prototype results are 100% synthetic and deterministic. Assumptions labeled throughout.",
        10, color=RGBColor(0x7C, 0x90, 0x9C), italic=True, space_after=0)


# ================================================================ SLIDE 2 — PROBLEM
s = slide()
kicker_title(s, "The Problem", "Oncology PA isn\u2019t paperwork. It interrupts active cancer therapy.", "02")

# left bullets
tb, tf = textbox(s, Inches(0.7), Inches(2.0), Inches(6.5), Inches(4.6))
bullets(tf, [
    ("AMA: ", "39 prior auths per physician per week, about 13 hours of staff time."),
    ("ASTRO: ", "92% of radiation oncologists report PA-driven treatment delays."),
    ("Delays are getting worse: ", "68% now last 5+ days, up from 52% in 2020."),
    ("ASTRO: ", "7% report PA contributed to a patient death."),
    ("JCO (oral oncolytics): ", "a new PA on an established drug raised discontinuation "
        "odds ~7.1x and added ~9.7 days to the next fill."),
], size=17.5, gap=14, line=1.1)

# bottom framing line
band = panel(s, Inches(0.7), Inches(5.95), Inches(6.5), Inches(0.85), fill=ACCENT_SOFT, line=None)
tb, tf = textbox(s, Inches(0.95), Inches(6.02), Inches(6.0), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "Risa already solves the initial PA. The gap is keeping approvals alive "
             "after they\u2019re granted.", 14.5, color=ACCENT_DK, bold=True, line=1.1, space_after=0)

# right stat stack
sx = Inches(7.65); sw = Inches(4.95)
stat_card(s, sx, Inches(2.0), sw, Inches(1.15), "39 / wk", "PAs per physician",
          "~13 staff-hours weekly (AMA)")
stat_card(s, sx, Inches(3.3), sw, Inches(1.15), "92%", "report treatment delays",
          "average 35% of patients affected (ASTRO)", accent=AMBER)
stat_card(s, sx, Inches(4.6), sw, Inches(1.15), "~7.1x", "discontinuation odds",
          "new PA on an established oral drug (JCO)", accent=RED)
footnote(s, "Sources: AMA 2024 PA survey · ASTRO 2024 PA survey · JCO 2010\u20132020 oral anticancer drugs. Full URLs on Sources slide.")


# ================================================================ SLIDE 3 — CURRENT SOLUTIONS
s = slide()
kicker_title(s, "Current Solutions + Shortcomings", "Everyone detects the work. Nobody executes the renewal.", "03")

# left column: what's used today + vendors
tb, tf = textbox(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(2.5))
add_para(tf, "What practices use today", 14, color=ACCENT_DK, bold=True, space_after=6)
bullets(tf, [
    "Manual auth-log spreadsheets",
    "PMS / RCM auth-tracking workqueues",
    "Eligibility & benefit modules",
    "EHR task and reminder queues",
    "Payer portals + CoverMyMeds status",
], size=14.5, gap=6, line=1.05)

tb, tf = textbox(s, Inches(0.7), Inches(4.55), Inches(6.0), Inches(2.2))
add_para(tf, "Adjacent vendors", 14, color=ACCENT_DK, bold=True, space_after=6)
add_para(tf, "Cohere · Rhyme · Myndshft · Waystar · Availity · Infinitus · "
             "Develop Health · Tennr · Co:Helm · Anterior",
         13.5, color=BODY, line=1.25, space_after=10)
add_para(tf, "They reduce friction in pieces of the workflow. None sit inside Risa\u2019s "
             "oncology PA loop with the EMR, approval record, form logic, and submission engine.",
         12.5, color=MUTED, italic=True, line=1.2, space_after=0)

# right: 2x2 matrix
mx, my = Inches(7.15), Inches(2.05)
mw, mh = Inches(5.4), Inches(4.05)
# axis labels
tb, tf = textbox(s, mx, my - Inches(0.0), mw, Inches(0.3))
# header row labels
cell_w = mw / 2
cell_h = mh / 2
# column headers
def axhdr(x, y, w, h, text, color=MUTED, size=11.5, bold=True, align=PP_ALIGN.CENTER):
    tb, tf = textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], text, size, color=color, bold=bold, align=align, space_after=0, line=1.0)

axhdr(mx, my - Inches(0.02), cell_w, Inches(0.32), "Detection only")
axhdr(mx + cell_w, my - Inches(0.02), cell_w, Inches(0.32), "Autonomous execution", color=ACCENT_DK)
# row labels (rotated-ish: just left side stacked)
axhdr(mx - Inches(1.35), my + Inches(0.34), Inches(1.3), cell_h, "Generic / cross-specialty", size=10.5, align=PP_ALIGN.RIGHT)
axhdr(mx - Inches(1.35), my + Inches(0.34) + cell_h, Inches(1.3), cell_h, "EMR-native oncology", size=10.5, align=PP_ALIGN.RIGHT)

grid_top = my + Inches(0.34)
cells = [
    (0, 0, "Spreadsheets, EHR reminders,\npayer portals", PANEL, MUTED, False),
    (1, 0, "Generic PA / RCM automation\n(Cohere, Waystar, Availity\u2026)", PANEL, MUTED, False),
    (0, 1, "Practice auth-log + manual\nstaff follow-up", PANEL, MUTED, False),
    (1, 1, "Auth Guardian\nEMR-native + client-approved execution", ACCENT, WHITE, True),
]
for col, row, txt, fill, txtcol, star in cells:
    cx = mx + col * cell_w
    cy = grid_top + row * cell_h
    c = rect(s, cx + Inches(0.04), cy + Inches(0.04), cell_w - Inches(0.08), cell_h - Inches(0.08),
             fill=fill, line=(ACCENT_DK if star else LINE), line_w=(2.0 if star else 1.0),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _round(c, 0.06)
    tb, tf = textbox(s, cx + Inches(0.2), cy + Inches(0.12), cell_w - Inches(0.36), cell_h - Inches(0.28),
                     anchor=MSO_ANCHOR.MIDDLE)
    lines = txt.split("\n")
    add_para(tf, lines[0], 13 if star else 11.5, color=txtcol, bold=star, align=PP_ALIGN.CENTER,
             space_after=2, line=1.05)
    if len(lines) > 1:
        add_para(tf, lines[1], 10.5 if star else 10, color=(RGBColor(0xDD,0xF0,0xF4) if star else MUTED),
                 align=PP_ALIGN.CENTER, space_after=0, line=1.05)

# matrix caption
tb, tf = textbox(s, mx, grid_top + mh + Inches(0.02), mw, Inches(0.5))
add_para(tf, "The empty quadrant is the wedge. Practices already have reminders \u2014 yet 5+ day "
             "oncology delays rose 52% \u2192 68%. The bottleneck is execution, not detection.",
         11, color=MUTED, italic=True, line=1.15, space_after=0)


# ================================================================ SLIDE 4 — MARKET SIZE
s = slide()
kicker_title(s, "Market Size", "Sized on monthly trigger cases per practice.", "04")

# Base case formula card (left)
panel(s, Inches(0.7), Inches(1.95), Inches(5.3), Inches(2.05))
tb, tf = textbox(s, Inches(0.95), Inches(2.08), Inches(4.85), Inches(1.85), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "BASE CASE \u00b7 ONE PRACTICE", 11, color=ACCENT_DK, bold=True, space_after=5)
add_para(tf, "750 active oral-oncolytic PA patients  \u00d7  12% monthly trigger  \u00d7  93% actionable",
         12.5, color=BODY, line=1.2, space_after=6)
add_para(tf, "\u2248 84 trigger cases / month", 19, color=ACCENT, bold=True, space_after=2)
tf.paragraphs[-1].runs[0].font.name = FONT_H
add_para(tf, "\u2248 $280K/month protected  ·  70 staff-hours/month",
         12.5, color=INK, bold=True, space_after=0)

# Risa book card
panel(s, Inches(0.7), Inches(4.15), Inches(5.3), Inches(1.95), fill=ACCENT_SOFT, line=None)
tb, tf = textbox(s, Inches(0.95), Inches(4.28), Inches(4.85), Inches(1.7), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "RISA BOOK \u00b7 NYCBS + ASTERA (placeholder)", 11, color=ACCENT_DK, bold=True, space_after=5)
add_para(tf, "1,950 combined patients \u00d7 12% \u00d7 93%", 12, color=BODY, space_after=6)
add_para(tf, "\u2248 218 triggers/mo  ·  $730K/mo protected  ·  $8.8M/yr",
         14.5, color=ACCENT_DK, bold=True, line=1.15, space_after=0)
tf.paragraphs[-1].runs[0].font.name = FONT_H

# Sensitivity table (right)
tb, tf = textbox(s, Inches(6.35), Inches(1.9), Inches(6.2), Inches(0.35))
add_para(tf, "Sensitivity \u00b7 per-practice monthly opportunity", 13, color=ACCENT_DK, bold=True, space_after=0)

rows = [
    ["Scenario", "Patients", "Trigger rate", "Cases/mo", "Lapses prevented", "Protected/mo"],
    ["Low", "400", "8%", "30", "10", "$100K"],
    ["Base", "750", "12%", "84", "28", "$280K"],
    ["High", "1,200", "16%", "179", "60", "$600K"],
]
tx, ty = Inches(6.35), Inches(2.35)
tw, th = Inches(6.25), Inches(1.85)
tbl_shape = s.shapes.add_table(len(rows), len(rows[0]), tx, ty, tw, th)
tbl = tbl_shape.table
# column widths
widths = [1.25, 1.0, 1.1, 0.95, 1.45, 1.1]
for i, wv in enumerate(widths):
    tbl.columns[i].width = Inches(wv)
for ri, row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.42)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc = cell.text_frame
        tfc.word_wrap = True
        p = tfc.paragraphs[0]
        p.text = val
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        is_base = (ri == 2)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(11 if ri else 10.5)
            r.font.bold = (ri == 0) or is_base
            if ri == 0:
                r.font.color.rgb = WHITE
            elif is_base:
                r.font.color.rgb = ACCENT_DK
            else:
                r.font.color.rgb = BODY
        # fill
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_DK
        elif is_base:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_SOFT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE

# TAM line
panel(s, Inches(6.35), Inches(4.45), Inches(6.25), Inches(1.6), fill=PANEL)
tb, tf = textbox(s, Inches(6.6), Inches(4.58), Inches(5.8), Inches(1.4), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "US TAM PROXY (modest, labeled assumptions)", 11, color=ACCENT_DK, bold=True, space_after=5)
add_para(tf, "1,000 practices \u00d7 84 cases/mo \u00d7 $100 per completed case",
         12, color=BODY, space_after=4)
add_para(tf, "\u2248 $101M / year", 18, color=ACCENT, bold=True, space_after=2)
tf.paragraphs[-1].runs[0].font.name = FONT_H
add_para(tf, "Revenue proxy is a sizing bridge, not a price. Range $50M\u2013$151M across 500\u20131,500 practices.",
         10, color=MUTED, italic=True, line=1.1, space_after=0)

footnote(s, "All figures rest on labeled assumptions: $10K avg claim, 33.6% reactive lapse rate, 2.5 hrs/lapse. NYCBS+Astera patient counts are placeholders.")


# ================================================================ SLIDE 5 — WHY RISA WINS + WHY NOW
s = slide()
kicker_title(s, "Why Risa Wins  +  Why Now", "Risa already owns the execution loop.", "05")

# left: why risa wins
tb, tf = textbox(s, Inches(0.7), Inches(1.95), Inches(6.4), Inches(0.4))
add_para(tf, "Why Risa wins", 16, color=ACCENT_DK, bold=True, space_after=0)
tb, tf = textbox(s, Inches(0.7), Inches(2.45), Inches(6.4), Inches(4.3))
bullets(tf, [
    ("Already inside the loop: ", "EMR data + PA records + CoverMyMeds + payer Q&A + write-back."),
    ("Zero new integration: ", "Auth Guardian adds trigger logic and an approval gate, not a new system."),
    ("Execution wedge: ", "others flag work; Risa files the renewal through its existing engine."),
    ("Additive revenue: ", "a new recurring workflow that doesn\u2019t cannibalize initial PA automation."),
    ("Fast to ship: ", "auth-expiry first, then coverage-expiry, then regimen-change."),
], size=15.5, gap=12, line=1.12)

# right: why now panel
panel(s, Inches(7.35), Inches(2.0), Inches(5.25), Inches(4.5), fill=ACCENT_SOFT, line=None)
tb, tf = textbox(s, Inches(7.65), Inches(2.2), Inches(4.7), Inches(4.1))
add_para(tf, "WHY NOW", 13, color=ACCENT_DK, bold=True, space_after=10)
add_para(tf, "CMS Interoperability & Prior Authorization Final Rule", 14.5, color=INK, bold=True,
         line=1.1, space_after=3)
add_para(tf, "CMS-0057-F, phased through 2026\u20132027. Impacted payers must support PA APIs and faster "
             "decisions \u2014 72 hours urgent, 7 days standard.",
         12.5, color=BODY, line=1.2, space_after=14)
add_para(tf, "State gold-card laws", 14.5, color=INK, bold=True, line=1.1, space_after=3)
add_para(tf, "Texas HB 3459 and the federal GOLD CARD Act push the same direction.",
         12.5, color=BODY, line=1.2, space_after=14)
add_para(tf, "Payer PA is modernizing \u2014 but practices still need provider-side systems that catch "
             "renewal risk early and execute before the gap.",
         12.5, color=ACCENT_DK, italic=True, bold=True, line=1.2, space_after=0)


# ================================================================ SLIDE 6 — THE SOLUTION
s = slide()
kicker_title(s, "The Solution", "Detect early. Hold for approval. File through the engine that already works.", "06")

# Three triggers row (left)
tb, tf = textbox(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(0.35))
add_para(tf, "Three triggers", 14, color=ACCENT_DK, bold=True, space_after=0)

trig = [
    ("Auth-expiry", "from Risa\u2019s own PA approval records", ACCENT),
    ("Insurance-expiry", "read from the EMR \u2014 not stored today", AMBER),
    ("Regimen-change", "monthly EMR re-check, lower priority", MUTED),
]
ty0 = Inches(2.3)
for i, (t, d, c) in enumerate(trig):
    yy = ty0 + Inches(0.78) * i
    card = rect(s, Inches(0.7), yy, Inches(5.7), Inches(0.66), fill=WHITE, line=LINE, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _round(card, 0.12)
    rect(s, Inches(0.7), yy, Inches(0.12), Inches(0.66), fill=c)
    tb, tf = textbox(s, Inches(1.0), yy + Inches(0.04), Inches(5.3), Inches(0.58), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, t, 14, color=INK, bold=True, space_after=1, line=1.0)
    add_para(tf, d, 11.5, color=MUTED, space_after=0, line=1.0)

# Approval gate emphasis
gate = panel(s, Inches(0.7), Inches(4.95), Inches(5.7), Inches(1.0), fill=ACCENT_SOFT, line=None)
tb, tf = textbox(s, Inches(0.95), Inches(5.02), Inches(5.25), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "CLIENT APPROVAL GATE", 12, color=ACCENT_DK, bold=True, space_after=2)
add_para(tf, "Nothing auto-files. Suppression holds stopped, transferred, or progressed cases for review.",
         12.5, color=BODY, line=1.15, space_after=0)

# Maria story (left bottom)
tb, tf = textbox(s, Inches(0.7), Inches(6.05), Inches(5.7), Inches(0.85))
add_para(tf, "Maria Alvarez \u00b7 Ibrance, Humana MA", 11.5, color=INK, bold=True, space_after=2)
add_para(tf, "Reactive: refiled at next visit \u2192 36-day gap + denial. "
             "Auth Guardian: fires 45 days early, client approves, filed on time \u2192 zero gap.",
         10.5, color=MUTED, line=1.15, space_after=0)

# Flow diagram (right column, vertical)
fx = Inches(6.85)
fw = Inches(5.7)
steps = [
    ("Approve PA", WHITE, INK, LINE),
    ("Store trigger dates", WHITE, INK, LINE),
    ("Daily date-watcher", WHITE, INK, LINE),
    ("Worklist highlights case", WHITE, INK, LINE),
    ("Client approval gate", ACCENT, WHITE, ACCENT_DK),
    ("Risa engine files renewal", WHITE, INK, LINE),
    ("EMR write-back + audit log", ACCENT_DK, WHITE, ACCENT_DK),
]
tb, tf = textbox(s, fx, Inches(1.9), fw, Inches(0.35))
add_para(tf, "How it flows", 14, color=ACCENT_DK, bold=True, space_after=0)

sy = Inches(2.3)
sh = Inches(0.52)
gap = Inches(0.135)
for i, (txt, fill, txtc, ln) in enumerate(steps):
    yy = sy + (sh + gap) * i
    c = rect(s, fx, yy, fw, sh, fill=fill, line=ln, line_w=(1.75 if fill in (ACCENT, ACCENT_DK) else 1.0),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _round(c, 0.18)
    tb, tf = textbox(s, fx + Inches(0.3), yy, fw - Inches(0.5), sh, anchor=MSO_ANCHOR.MIDDLE)
    bold = fill in (ACCENT, ACCENT_DK)
    add_para(tf, ("\u2192  " if not bold else "\u2713  ") + txt if bold else txt,
             13.5, color=txtc, bold=bold, space_after=0, line=1.0)
    # connector chevron
    if i < len(steps) - 1:
        ar = rect(s, fx + fw/2 - Inches(0.09), yy + sh + Inches(0.0), Inches(0.18), gap,
                  fill=None, shape=MSO_SHAPE.DOWN_ARROW)
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xB6, 0xC4, 0xCE)
        ar.line.fill.background()


# ================================================================ SLIDE 7 — IMPACT
s = slide()
kicker_title(s, "Impact \u00b7 The Backtest", "Reactive vs Auth Guardian on the same 150-case cohort.", "07")

# stat cards row
cy = Inches(1.95)
cw = Inches(2.86); ch = Inches(1.25); cgap = Inches(0.13)
cx0 = Inches(0.7)
cards = [
    ("47 \u2192 0", "Lapsed authorizations", "33.6% \u2192 0%", RED),
    ("2,024 \u2192 0", "Therapy-gap days", "eliminated", AMBER),
    ("$470K", "Claim value protected", "@ $10K/claim (assumption)", ACCENT),
    ("118 hrs", "Staff scramble saved", "@ 2.5 hrs/lapse", GREEN),
]
for i, (big, lab, sub, c) in enumerate(cards):
    stat_card(s, cx0 + (cw + cgap) * i, cy, cw, ch, big, lab, sub, accent=c)

# bar chart 47 -> 0 (left lower)
chart_data = CategoryChartData()
chart_data.categories = ["Reactive (today)", "Auth Guardian"]
chart_data.add_series("Lapsed authorizations", (47, 0))
gx, gy = Inches(0.7), Inches(3.55)
gw, gh = Inches(5.6), Inches(3.0)
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, gx, gy, gw, gh, chart_data)
chart = gframe.chart
chart.has_legend = False
chart.has_title = True
chart.chart_title.text_frame.text = "Lapsed authorizations: 47 \u2192 0"
ctp = chart.chart_title.text_frame.paragraphs[0]
ctp.runs[0].font.size = Pt(13); ctp.runs[0].font.bold = True
ctp.runs[0].font.color.rgb = INK; ctp.runs[0].font.name = FONT
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(16)
plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = INK
plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
plot.gap_width = 120
# color points individually
series = chart.series[0]
for idx, col in enumerate([RED, ACCENT]):
    pt = series.points[idx]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = col
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(11); cat_ax.tick_labels.font.color.rgb = BODY
cat_ax.has_major_gridlines = False
val_ax = chart.value_axis
val_ax.has_major_gridlines = False
val_ax.visible = False
val_ax.maximum_scale = 55

# right: deltas list
panel(s, Inches(6.55), Inches(3.55), Inches(6.05), Inches(3.0))
tb, tf = textbox(s, Inches(6.85), Inches(3.72), Inches(5.5), Inches(2.7), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "WHAT THE PROTOTYPE PROVES", 11.5, color=ACCENT_DK, bold=True, space_after=8)
bullets(tf, [
    ("47 no-auth denials ", "avoided (now 0)."),
    ("10 / 10 suppression triage ", "\u2014 stopped/transferred/progressed cases held for review."),
    ("140 client-approval touches ", "\u2014 nothing auto-files."),
    ("Synthetic, deterministic prototype: ", "numbers prove the mechanism, not a market forecast."),
], size=13, gap=10, line=1.15)

footnote(s, "Prototype data is 100% synthetic and deterministic (seed 20260621). Dollar and hour figures use labeled assumptions: $10K/claim, 2.5 hrs/lapse.")


# ================================================================ SLIDE 8 — VISION
s = slide()
kicker_title(s, "Vision \u00b7 What\u2019s Next", "From continuity layer to revenue-integrity layer.", "08")

tb, tf = textbox(s, Inches(0.7), Inches(2.1), Inches(7.4), Inches(4.2))
bullets(tf, [
    ("Approval is not payment. ", "The same proactive engine extends to claim reconciliation."),
    ("Add a PMS feed ", "and Auth Guardian watches paid vs denied, not just approved vs lapsed."),
    ("Continuity layer \u2192 revenue-integrity layer ", "across the full post-approval lifecycle."),
    ("Same wedge: ", "detect early, hold for client approval, then execute through Risa\u2019s engine."),
], size=18, gap=18, line=1.15)

# right vertical progression
fx = Inches(8.6); fw = Inches(4.0)
stages = [
    ("Today", "Initial PA automation (BOSS)", MUTED),
    ("Auth Guardian", "Proactive renewal continuity", ACCENT),
    ("Next", "Claim reconciliation on PMS feed", ACCENT_DK),
    ("Vision", "Revenue-integrity layer", INK),
]
sy = Inches(2.0)
for i, (t, d, c) in enumerate(stages):
    yy = sy + Inches(1.18) * i
    card = rect(s, fx, yy, fw, Inches(0.92), fill=WHITE, line=c, line_w=1.75,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _round(card, 0.1)
    rect(s, fx, yy, Inches(0.13), Inches(0.92), fill=c)
    tb, tf = textbox(s, fx + Inches(0.35), yy + Inches(0.08), fw - Inches(0.5), Inches(0.78),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, t, 13, color=c, bold=True, space_after=2, line=1.0)
    add_para(tf, d, 12.5, color=BODY, space_after=0, line=1.05)
    if i < len(stages) - 1:
        ar = rect(s, fx + fw/2 - Inches(0.09), yy + Inches(0.93), Inches(0.18), Inches(0.24),
                  shape=MSO_SHAPE.DOWN_ARROW)
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xB6, 0xC4, 0xCE)
        ar.line.fill.background()


# ================================================================ SLIDE 9 — SOURCES
s = slide()
rect(s, Inches(0), Inches(0), Inches(0.18), EMU_H, fill=ACCENT)
tb, tf = textbox(s, Inches(0.7), Inches(0.5), Inches(11.4), Inches(0.9))
add_para(tf, "SOURCES & METHODOLOGY", 13, color=ACCENT, bold=True, space_after=4)
add_para(tf, "Sources", 28, color=INK, bold=True, line=1.0)

sources = [
    "1. AMA \u2014 Fixing prior auth: nearly 40 prior authorizations a week",
    "   ama-assn.org/practice-management/prior-authorization/fixing-prior-auth-nearly-40-prior-authorizations-week-way",
    "2. AMA \u2014 Prior Authorization Reform Progress Update (PDF)",
    "3. ASTRO \u2014 Prior Authorization Survey 2024 Executive Summary (PDF)",
    "4. AJMC \u2014 Prior Authorization Delays Cause Serious Harm to Patients With Cancer",
    "5. The ASCO Post \u2014 Prior Authorization Delays May Lead to Severe Consequences",
    "6. PubMed/JCO \u2014 Association of Prior Authorization and Price With Access to Oral Anticancer Drugs (38086013)",
    "7. Go-Flow \u2014 Prior Authorization Denial Reasons: 2026 Data + Fixes",
    "8. The SORSO \u2014 Why Are Claims Denied?",
    "9. Nirmitee \u2014 Healthcare Denial Root Cause Analysis",
    "10. Klivira \u2014 Oncology Infusion Prior Authorization Playbook: Regimen to Re-Auth",
    "11. PharmaDossier \u2014 Renewal Denial: Stable Disease, Continued Therapy",
    "12. CMS \u2014 Interoperability and Prior Authorization Final Rule (CMS-0057-F)",
    "13. Texas Dept. of Insurance \u2014 Gold Carding Prior Authorization Exemptions",
    "14. AMA \u2014 Gold Card Act",
]
tb, tf = textbox(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(4.6))
for i, src in enumerate(sources):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    is_url = src.strip().startswith(("ama-assn", "covermymeds")) or "/" in src and src.startswith("   ")
    setpara(p, src, 11 if not src.startswith("   ") else 9.5,
            color=(MUTED if src.startswith("   ") else BODY),
            italic=src.startswith("   "), space_after=(2 if src.startswith("   ") else 5), line=1.05)

panel(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.85), fill=PANEL)
tb, tf = textbox(s, Inches(0.95), Inches(6.33), Inches(11.4), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
add_para(tf, "Prototype data is 100% synthetic and deterministic (mulberry32 seed 20260621, watcher date 2026-06-21). "
             "All dollar, hour, patient-count, and TAM figures use clearly labeled assumptions.",
         11, color=MUTED, italic=True, line=1.2, space_after=0)


# ---------------------------------------------------------------- save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "auth-guardian.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
